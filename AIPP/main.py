import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
import json
import imageGeneration
import chatGPT
#import gui
import os

token = ""
prompt = ""
downloadsFolder = os.path.join(os.path.expanduser('~'), 'Downloads')

def getUserData(userToken,userPrompt):
    global token,prompt

    token = userToken
    prompt = userPrompt

    print(token)
    print(prompt)

    createPresentation()


def createPresentation():

    chatGPT.callChat(prompt)

    # load the data from the json file
    with open('response.json', 'r') as data:
        slideData = json.load(data)

    # Initialize Variables
    prs = Presentation("preset1.pptx")

    allBulletPoints = ""
    slideIndex = 0

    # Loop through the data and create each slide with its corresponding data
    for currentSlide in slideData:
        #print(slideIndex)
        if slideIndex == 0:
            slide_layout = prs.slide_layouts[9]
            slide = prs.slides.add_slide(slide_layout)

            title = slide.shapes.title
            subtitle = slide.placeholders[2]

            title.text = currentSlide["title"]
            subtitle.text = currentSlide["subtitle"]

            # generate images and add them to the presentation
            imageGeneration.generateImage(currentSlide["image"],token)

            img_path = currentSlide["image"].strip()[0] + ".png"
            photo = slide.placeholders[1]

            photo.insert_picture(img_path)
            #pic = slide.shapes.add_picture(img_path, left, top, width, height)


            slideIndex = slideIndex + 1
        else:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            print(prs.slide_layouts[1].name)

            # Set the slide title, subtitle and bullet points
            title = slide.shapes.title
            title.text = currentSlide["title"]

            # Subtitle
            subtitle = None
            for shape in slide.shapes:
                print(shape.name)
                print(shape.shape_id)
                if shape.name == "Text Placeholder 3":
                    subtitle = shape
                    break
            subtitle.text = currentSlide["subtitle"]

            for paragraph in subtitle.text_frame.paragraphs:
                paragraph.font.size = Pt(20)


            # BulletPoints
            bulletpoints = slide.placeholders[1]


            # Add all the bullet-points in one string.
            for bullet in currentSlide["bullets"]:
                print(bullet)
                allBulletPoints = allBulletPoints + bullet + "\n"

            # Add the string to the bullet-point thing in the presentation
            bulletpoints.text = allBulletPoints
            for paragraph in bulletpoints.text_frame.paragraphs:
                paragraph.font.size = Pt(15)

            photo = None
            for shape in slide.shapes:
                print(shape.name)
                print(shape.shape_id)
                if shape.name == "Picture Placeholder 4":
                    photo = shape
                    break

            #generate images and add them to the presentation
            imageGeneration.generateImage(currentSlide["image"],token)

            img_path = currentSlide["image"].strip()[0] + ".png"
            photo.insert_picture(img_path)
            #photo = slide.placeholders[1]

            '''left = top = Inches(1)
            pic = slide.shapes.add_picture(img_path, left, top)
    
            pic.left = Inches(2)
            pic.top = Inches(3)
            pic.width = Inches(4)
            pic.height = Inches(4)'''
            '''photo = slide.placeholders[1]
    
            photo.insert_picture(img_path)'''
            # Change slide index and clear the string, so it doesn't stack up all the bullet-points.
            slideIndex = slideIndex + 1
            allBulletPoints = ""

    #prs.save('test7.pptx')
    prs.save(os.path.join(downloadsFolder, 'finalPresentation.pptx'))
    print("done")

    jsonFile = "response.json"

    with open(jsonFile, 'w') as f:
        pass